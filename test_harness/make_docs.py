# -*- coding: utf-8 -*-
"""Genera documentos de prueba 'difíciles' para evaluar la conversión a Markdown."""
import io, os, datetime
from pathlib import Path

OUT = Path(__file__).parent / "docs"
OUT.mkdir(exist_ok=True)

# ── 1. PDF con texto, headings, tabla y lista ────────────────────────────────
import fitz

doc = fitz.open()
page = doc.new_page()
y = 60
page.insert_text((50, y), "Informe Anual 2025", fontsize=22, fontname="helv")
y += 40
page.insert_text((50, y), "1. Resumen Ejecutivo", fontsize=16, fontname="hebo")
y += 28
texto = ("La compañía alcanzó ingresos por USD 4,5 millones durante el ejercicio 2025, "
         "un crecimiento del 23% respecto al año anterior. El margen EBITDA se situó en 31%.")
for line in [texto[i:i+90] for i in range(0, len(texto), 90)]:
    page.insert_text((50, y), line, fontsize=11); y += 16
y += 12
page.insert_text((50, y), "2. Resultados por Segmento", fontsize=16, fontname="hebo")
y += 28
# tabla dibujada con líneas
rows = [["Segmento", "Ingresos (USD M)", "Margen"],
        ["Retail", "2,1", "28%"],
        ["Corporativo", "1,7", "35%"],
        ["Digital", "0,7", "41%"]]
x0, col_w, row_h = 50, 140, 22
for ri, row in enumerate(rows):
    for ci, cell in enumerate(row):
        page.insert_text((x0 + ci*col_w + 6, y + ri*row_h + 15), cell, fontsize=10,
                         fontname="hebo" if ri == 0 else "helv")
for ri in range(len(rows)+1):
    page.draw_line((x0, y + ri*row_h), (x0 + 3*col_w, y + ri*row_h))
for ci in range(4):
    page.draw_line((x0 + ci*col_w, y), (x0 + ci*col_w, y + len(rows)*row_h))
y += len(rows)*row_h + 30
page.insert_text((50, y), "3. Prioridades 2026", fontsize=16, fontname="hebo")
y += 26
for item in ["• Expansión a mercados andinos", "• Migración cloud de sistemas core",
             "• Programa de eficiencia: meta -8% en gastos"]:
    page.insert_text((62, y), item, fontsize=11); y += 18

# página 2 con texto y caracteres especiales
p2 = doc.new_page()
p2.insert_text((50, 60), "Anexo A: Notas metodológicas", fontsize=16, fontname="hebo")
p2.insert_text((50, 95), "Cifras en pesos chilenos (CLP$). Tipo de cambio: 945 CLP/USD.", fontsize=11)
p2.insert_text((50, 115), "Símbolos de prueba: 100% | a*b | c_d | #tag | <html> & \"quotes\"", fontsize=11)
doc.save(OUT / "informe_texto.pdf")
doc.close()

# ── 2. PDF "escaneado" (solo imagen, sin capa de texto) ──────────────────────
src = fitz.open(OUT / "informe_texto.pdf")
scan = fitz.open()
for p in src:
    pix = p.get_pixmap(dpi=120)
    np_ = scan.new_page(width=p.rect.width, height=p.rect.height)
    np_.insert_image(np_.rect, pixmap=pix)
scan.save(OUT / "informe_escaneado.pdf")
scan.close(); src.close()

# ── 2b. PDF con portada-imagen y resto texto (caso borde para detección OCR) ─
mix = fitz.open()
cover = mix.new_page()
pix = fitz.open(OUT / "informe_escaneado.pdf")[0].get_pixmap(dpi=100)
cover.insert_image(cover.rect, pixmap=pix)
tp = mix.new_page()
tp.insert_text((50, 60), "Capítulo 1: Contenido real en texto", fontsize=16, fontname="hebo")
tp.insert_text((50, 95), "Este PDF tiene portada escaneada pero el resto es texto digital.", fontsize=11)
for i in range(20):
    tp.insert_text((50, 120 + i*15), f"Línea de contenido número {i+1} con texto digital extraíble.", fontsize=10)
mix.save(OUT / "portada_imagen_resto_texto.pdf")
mix.close()

# ── 3. DOCX con estilos, tabla, lista, imagen ─────────────────────────────────
from docx import Document
from docx.shared import Inches

d = Document()
d.add_heading("Política de Vacaciones", level=1)
d.add_paragraph("Vigente desde el 1 de marzo de 2026. Aplica a todos los colaboradores con contrato indefinido.")
d.add_heading("Días disponibles", level=2)
t = d.add_table(rows=4, cols=3)
t.style = "Table Grid"
data = [["Antigüedad", "Días hábiles", "Observación"],
        ["< 1 año", "proporcional", "1,25 días/mes"],
        ["1–5 años", "15", "estándar legal"],
        ["> 5 años", "15 + 1 por trienio", "progresivas | tope 20"]]
for ri, row in enumerate(data):
    for ci, val in enumerate(row):
        t.rows[ri].cells[ci].text = val
d.add_heading("Procedimiento", level=2)
for s in ["Solicitar con 15 días de anticipación vía portal RRHH.",
          "Aprobación del jefe directo (máx. 3 días hábiles).",
          "Confirmación automática por correo."]:
    d.add_paragraph(s, style="List Number")
d.add_heading("Excepciones", level=2)
for s in ["Periodos de cierre contable (enero y julio).",
          "Equipos con dotación crítica < 70%."]:
    d.add_paragraph(s, style="List Bullet")
p = d.add_paragraph()
run = p.add_run()
# imagen pequeña: cuadrado PNG
img_doc = fitz.open(); ip = img_doc.new_page(width=200, height=120)
ip.draw_rect(fitz.Rect(10, 10, 190, 110), color=(0, 0.35, 0.44), fill=(0.06, 0.71, 0.75))
ip.insert_text((40, 65), "ORGANIGRAMA", fontsize=14, color=(1, 1, 1))
png_bytes = ip.get_pixmap(dpi=96).tobytes("png")
img_doc.close()
run.add_picture(io.BytesIO(png_bytes), width=Inches(2.5))
caption = d.add_paragraph("Figura 1: organigrama del área.")
d.add_paragraph("Texto con formato: ")
p2 = d.add_paragraph()
p2.add_run("negrita").bold = True
p2.add_run(" y ")
p2.add_run("cursiva").italic = True
p2.add_run(" y símbolos: 50% | a_b | #ref.")
d.save(OUT / "politica_vacaciones.docx")

# ── 4. PPTX con títulos, bullets, tabla, imagen y notas ──────────────────────
from pptx import Presentation
from pptx.util import Inches as PInches, Pt

prs = Presentation()
# Slide 1: título
s1 = prs.slides.add_slide(prs.slide_layouts[0])
s1.shapes.title.text = "Plan Comercial 2026"
s1.placeholders[1].text = "Gerencia Comercial — Enero 2026"
# Slide 2: bullets
s2 = prs.slides.add_slide(prs.slide_layouts[1])
s2.shapes.title.text = "Objetivos del Año"
tf = s2.placeholders[1].text_frame
tf.text = "Crecer 18% en ventas netas"
for txt, lvl in [("Nuevos clientes: +120 cuentas", 0),
                 ("Segmento pyme: +80", 1),
                 ("Segmento corporativo: +40", 1),
                 ("Churn bajo 5%", 0)]:
    para = tf.add_paragraph(); para.text = txt; para.level = lvl
# Slide 3: tabla
s3 = prs.slides.add_slide(prs.slide_layouts[5])
s3.shapes.title.text = "Metas por Trimestre"
tbl = s3.shapes.add_table(5, 3, PInches(1), PInches(1.8), PInches(8), PInches(3)).table
hdr = [["Trimestre", "Meta (UF)", "Responsable"],
       ["Q1", "12.500", "A. Rojas"],
       ["Q2", "14.000", "A. Rojas"],
       ["Q3", "15.800", "M. Pérez"],
       ["Q4", "18.200", "M. Pérez"]]
for ri, row in enumerate(hdr):
    for ci, val in enumerate(row):
        tbl.cell(ri, ci).text = val
# Slide 4: imagen + notas
s4 = prs.slides.add_slide(prs.slide_layouts[5])
s4.shapes.title.text = "Mapa de Cobertura"
s4.shapes.add_picture(io.BytesIO(png_bytes), PInches(2), PInches(2), width=PInches(4))
s4.notes_slide.notes_text_frame.text = "Mencionar que la cobertura sur se licita en marzo."
prs.save(OUT / "plan_comercial.pptx")

# ── 5. XLSX con 2 hojas, fechas, floats, pipes y saltos de línea ──────────────
from openpyxl import Workbook

wb = Workbook()
ws = wb.active
ws.title = "Dotación"
ws.append(["Nombre", "Área", "Ingreso", "Sueldo (CLP)", "Notas"])
ws.append(["María Paz Soto", "Finanzas", datetime.date(2021, 3, 15), 2450000, "jornada completa"])
ws.append(["Jorge Díaz", "TI | Infraestructura", datetime.date(2023, 11, 2), 3100000.0, "linea1\nlinea2"])
ws.append(["Antonia Núñez", "Comercial", datetime.datetime(2025, 6, 1, 9, 30), 1980000.5, None])
ws2 = wb.create_sheet("Resumen")
ws2.append(["Indicador", "Valor"])
ws2.append(["Dotación total", 3])
ws2.append(["Costo mensual", "=SUM(Dotación!D2:D4)"])  # fórmula
ws2.append([None, None])  # fila vacía intermedia
ws2.append(["Actualizado", datetime.date(2026, 6, 1)])
ws3 = wb.create_sheet("Vacía")
wb.save(OUT / "dotacion.xlsx")

# ── 6. TXT en latin-1 (acentos) — caso típico de exportes legacy ─────────────
txt = ("Acta de reunión — Comité de Personas\n"
       "Asistentes: Sebastián Ibáñez, Ma. José Peña\n"
       "Acuerdos:\n"
       "1) Aprobación del plan de capacitación.\n"
       "2) Revisión de banda salarial en práctica.\n")
(OUT / "acta_utf8.txt").write_bytes(txt.encode("utf-8"))
(OUT / "acta_latin1.txt").write_bytes(txt.encode("cp1252"))

print("Documentos generados:")
for f in sorted(OUT.iterdir()):
    print(f"  {f.name}  ({f.stat().st_size:,} bytes)")
